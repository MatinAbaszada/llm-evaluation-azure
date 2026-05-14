"""
Build the Bachelor Thesis defense presentation as a .pptx file.

Generates: Presentation/Thesis_Defense.pptx

Includes:
  - 15 main slides + 5 backup slides
  - Azure-themed color palette
  - Charts from /charts folder
  - Title-bar / callout / table styling
  - Constructor University logo extracted from Thesis Guideline.pptx
  - Entrance animations (fade-in on click) — shapes that form a single
    visual element (card, callout, diagram node) appear together as
    one click step.
"""

from __future__ import annotations
import math
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ------------------------------------------------------------------
# Paths
# ------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent
PRES_DIR = ROOT / "Presentation"
CHARTS = ROOT / "charts"
GUIDELINE_PPTX = PRES_DIR / "Thesis Guideline.pptx"
LOGO_PATH = PRES_DIR / "cu_logo_extracted.png"
OUT_PATH = PRES_DIR / "Thesis_Defense.pptx"

# ------------------------------------------------------------------
# Color palette
# ------------------------------------------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NAVY        = RGBColor(0x0B, 0x1F, 0x4B)
AZURE       = RGBColor(0x00, 0x78, 0xD4)
NEAR_BLACK  = RGBColor(0x21, 0x21, 0x21)
ICE_BLUE    = RGBColor(0xEB, 0xF4, 0xFF)
GREEN       = RGBColor(0x10, 0x7C, 0x10)
ORANGE_RED  = RGBColor(0xD8, 0x3B, 0x01)
AMBER       = RGBColor(0xFF, 0xB9, 0x00)
GRAY        = RGBColor(0x60, 0x5E, 0x5C)
AMBER_FILL  = RGBColor(0xFF, 0xF4, 0xCE)
GREEN_FILL  = RGBColor(0xE6, 0xF4, 0xEA)
RED_FILL    = RGBColor(0xFD, 0xE7, 0xE9)
ALT_ROW     = RGBColor(0xF7, 0xF9, 0xFC)

FONT = "Segoe UI"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.95)


# ------------------------------------------------------------------
# Animation grouping registry
# ------------------------------------------------------------------
# Maps id(slide) -> list of groups; each group = list of shape objects
# that should animate (fade in) together on a single click.
_ANIM_GROUPS: dict[int, list[list]] = {}
# Shapes registered as "always visible" (no animation, e.g. background,
# header bar, footer text).
_ALWAYS_VISIBLE: dict[int, set] = {}


def _gkey(slide):
    return id(slide)


def register_group(slide, *shapes):
    """Mark `shapes` as a single animation group (one click reveals them)."""
    if not shapes:
        return
    _ANIM_GROUPS.setdefault(_gkey(slide), []).append(
        [s for s in shapes if s is not None])


def register_always_visible(slide, *shapes):
    """Mark shapes as always visible (no animation, present from slide start).
    Stores stable XML shape_ids — Python id() of pptx wrappers is unstable
    because slide.shapes returns fresh wrapper objects on each iteration."""
    s_set = _ALWAYS_VISIBLE.setdefault(_gkey(slide), set())
    for sh in shapes:
        if sh is None:
            continue
        try:
            s_set.add(int(sh.shape_id))
        except Exception:
            pass


# ------------------------------------------------------------------
# Logo extraction
# ------------------------------------------------------------------
def extract_logo():
    if LOGO_PATH.exists() and LOGO_PATH.stat().st_size > 1000:
        return LOGO_PATH
    if not GUIDELINE_PPTX.exists():
        print(f"[warn] Guideline pptx not found: {GUIDELINE_PPTX}")
        return None

    # Fallback 1: python-pptx
    try:
        prs = Presentation(str(GUIDELINE_PPTX))
        for slide in prs.slides[:3]:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        LOGO_PATH.write_bytes(shape.image.blob)
                        print(f"[info] Logo extracted -> {LOGO_PATH}")
                        return LOGO_PATH
                    except Exception:
                        pass
    except Exception as e:
        print(f"[warn] python-pptx logo extraction failed: {e}")

    # Fallback 2: largest image in ppt/media
    try:
        with zipfile.ZipFile(str(GUIDELINE_PPTX)) as z:
            media = [n for n in z.namelist()
                     if n.startswith("ppt/media/")
                     and n.lower().endswith((".png", ".jpg", ".jpeg"))]
            candidates = sorted(((z.getinfo(n).file_size, n) for n in media),
                                reverse=True)
            for size, name in candidates:
                LOGO_PATH.write_bytes(z.read(name))
                print(f"[info] Logo extracted from zip -> {LOGO_PATH} "
                      f"({name}, {size} bytes)")
                return LOGO_PATH
    except Exception as e:
        print(f"[warn] zip logo extraction failed: {e}")
    return None


# ------------------------------------------------------------------
# Text / shape primitives
# ------------------------------------------------------------------
def set_run(run, text, *, size=20, bold=False, italic=False,
            color=NEAR_BLACK, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, size=20, bold=False,
                italic=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, italic=italic,
            color=color, font=font)
    return tb


def add_filled_rect(slide, left, top, width, height, fill_color,
                    line_color=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        if line_width is not None:
            shp.line.width = line_width
    shp.shadow.inherit = False
    return shp


def add_white_background(slide):
    bg = add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    register_always_visible(slide, bg)
    return bg


def add_header(slide, title_text, slide_num=None, total=16, section_label=None):
    header_shapes = []
    bar = add_filled_rect(slide, 0, 0, SLIDE_W, HEADER_H, NAVY)
    accent = add_filled_rect(slide, 0, HEADER_H, SLIDE_W, Inches(0.06), AZURE)
    title_tb = add_textbox(slide, Inches(0.5), Inches(0.15), Inches(11),
                           Inches(0.7), title_text, size=26, bold=True,
                           color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    header_shapes.extend([bar, accent, title_tb])

    if slide_num is not None:
        num_tb = add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.7),
                             Inches(0.35), f"{slide_num} / {total}", size=11,
                             color=GRAY, align=PP_ALIGN.RIGHT, italic=True)
        header_shapes.append(num_tb)

    if section_label:
        sec_tb = add_textbox(slide, Inches(0.5), Inches(7.05), Inches(6),
                             Inches(0.35), section_label, size=11, color=GRAY,
                             italic=True)
        header_shapes.append(sec_tb)

    register_always_visible(slide, *header_shapes)


def add_callout(slide, left, top, width, height, text, *,
                fill=ICE_BLUE, border=AZURE, text_color=NEAR_BLACK,
                bold=False, size=16, align=PP_ALIGN.LEFT):
    """Returns the list of constituent shapes (box + accent bar + text)."""
    box = add_filled_rect(slide, left, top, width, height, fill,
                          line_color=border, line_width=Pt(0.75))
    bar = add_filled_rect(slide, left, top, Inches(0.08), height, border)
    tb = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.05),
                                  width - Inches(0.25), height - Inches(0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    lines = text.split("\n")
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align
        set_run(p.add_run(), line, size=size, bold=bold, color=text_color)
    return [box, bar, tb]


def add_bullet_list(slide, left, top, width, height, items, *,
                    size=18, color=NEAR_BLACK, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(8)
        r1 = p.add_run()
        set_run(r1, "• ", size=size, bold=True, color=AZURE)
        r2 = p.add_run()
        set_run(r2, item, size=size, color=color)
    return tb


def add_picture_fitted(slide, image_path, left, top, max_w, max_h):
    from PIL import Image
    with Image.open(image_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        w, h = max_w, int(max_w / aspect)
    else:
        h, w = max_h, int(max_h * aspect)
    cx = left + (max_w - w) // 2
    cy = top + (max_h - h) // 2
    return slide.shapes.add_picture(str(image_path), cx, cy, width=w, height=h)


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    s = prs.slides.add_slide(layout)
    add_white_background(s)
    return s


# ------------------------------------------------------------------
# Compound widgets (return their constituent shapes for grouping)
# ------------------------------------------------------------------
def add_card_block(slide, x, y, w, h, *, fill=ICE_BLUE, border=AZURE,
                   border_w=Pt(0.75), badge_num=None, title=None,
                   title_color=NAVY, body=None, body_size=14):
    shapes = []
    bg = add_filled_rect(slide, x, y, w, h, fill, line_color=border,
                         line_width=border_w)
    shapes.append(bg)
    title_x = x + Inches(0.25)
    title_y = y + Inches(0.2)
    if badge_num is not None:
        badge = add_filled_rect(slide, x + Inches(0.2), y + Inches(0.2),
                                Inches(0.45), Inches(0.45), border)
        bt = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2),
                                      Inches(0.45), Inches(0.45))
        bt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = bt.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), str(badge_num), size=18, bold=True, color=WHITE)
        shapes.extend([badge, bt])
        title_x = x + Inches(0.8)
    if title:
        t = add_textbox(slide, title_x, title_y,
                        w - (title_x - x) - Inches(0.2),
                        Inches(0.5), title, size=18, bold=True,
                        color=title_color, anchor=MSO_ANCHOR.MIDDLE)
        shapes.append(t)
    if body:
        body_y = y + Inches(0.85) if title else y + Inches(0.25)
        b = add_textbox(slide, x + Inches(0.25), body_y,
                        w - Inches(0.5), h - (body_y - y) - Inches(0.15),
                        body, size=body_size, color=NEAR_BLACK)
        shapes.append(b)
    return shapes


def add_diag_node(slide, x, y, w, h, text, *, color=AZURE, fill=ICE_BLUE,
                  size=13, bold=True, body_color=NEAR_BLACK):
    bg = add_filled_rect(slide, x, y, w, h, fill,
                         line_color=color, line_width=Pt(1.0))
    t = add_textbox(slide, x, y, w, h, text, size=size, bold=bold,
                    color=body_color, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    return [bg, t]


def add_arrow_down(slide, x, y, w=Inches(0.3), h=Inches(0.2), color=AZURE):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


# ------------------------------------------------------------------
# Animation timing XML injection
# ------------------------------------------------------------------
ANIM_NS = (
    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
)


def _build_click_step(group_idx, shape_ids):
    """One click step = first shape is clickEffect, rest are withEffect (so
    they animate simultaneously)."""
    base = group_idx * 1000 + 100
    sub_effects = []
    for i, sid in enumerate(shape_ids):
        node_type = "clickEffect" if i == 0 else "withEffect"
        eid = base + 10 + i * 6
        sub_effects.append(f'''
          <p:par>
            <p:cTn id="{eid}" presetID="10" presetClass="entr" presetSubtype="0"
                   fill="hold" grpId="0" nodeType="{node_type}">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
              <p:childTnLst>
                <p:set>
                  <p:cBhvr>
                    <p:cTn id="{eid+1}" dur="1" fill="hold">
                      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    </p:cTn>
                    <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                    <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                  </p:cBhvr>
                  <p:to><p:strVal val="visible"/></p:to>
                </p:set>
                <p:anim calcmode="lin" valueType="num">
                  <p:cBhvr additive="base">
                    <p:cTn id="{eid+2}" dur="500" fill="hold"/>
                    <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                    <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>
                  </p:cBhvr>
                  <p:tavLst>
                    <p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>
                    <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav>
                  </p:tavLst>
                </p:anim>
              </p:childTnLst>
            </p:cTn>
          </p:par>''')

    return f'''
      <p:par>
        <p:cTn id="{base}" fill="hold">
          <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
          <p:childTnLst>
            <p:par>
              <p:cTn id="{base+1}" fill="hold">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>{''.join(sub_effects)}</p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
      </p:par>'''


def add_grouped_animations(slide):
    groups = _ANIM_GROUPS.get(_gkey(slide), [])
    always_xml_ids: set[int] = set(_ALWAYS_VISIBLE.get(_gkey(slide), set()))

    # Resolve grouped shapes to stable XML shape_ids (used for dedup only)
    grouped_xml_ids: set[int] = set()
    for g in groups:
        for sh in g:
            try:
                grouped_xml_ids.add(int(sh.shape_id))
            except Exception:
                pass

    # Fallback: any shape on the slide that isn't already registered in a
    # group AND isn't marked always-visible becomes its own click step.
    final_groups = list(groups)
    for shape in slide.shapes:
        try:
            sid = int(shape.shape_id)
        except Exception:
            continue
        if sid in grouped_xml_ids or sid in always_xml_ids:
            continue
        final_groups.append([shape])

    if not final_groups:
        return

    click_steps = []
    bld_entries = []
    for gi, group in enumerate(final_groups):
        sids = []
        for sh in group:
            try:
                sids.append(int(sh.shape_id))
            except Exception:
                continue
        if not sids:
            continue
        click_steps.append(_build_click_step(gi, sids))
        for sid in sids:
            bld_entries.append(f'<p:bldP spid="{sid}" grpId="0"/>')

    if not click_steps:
        return

    timing_xml = f'''<p:timing {ANIM_NS}>
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    {''.join(click_steps)}
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
                <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
      <p:bldLst>
        {''.join(bld_entries)}
      </p:bldLst>
    </p:timing>'''

    timing_el = etree.fromstring(timing_xml)
    sld = slide._element
    for old in sld.findall(qn('p:timing')):
        sld.remove(old)
    sld.append(timing_el)


# ------------------------------------------------------------------
# Slide builders
# ------------------------------------------------------------------
def slide_title(prs, logo_path):
    s = blank_slide(prs)
    navy_bg = add_filled_rect(s, 0, 0, SLIDE_W, Inches(4.5), NAVY)
    rule = add_filled_rect(s, Inches(2), Inches(4.45), Inches(9.333),
                           Inches(0.05), AZURE)
    register_always_visible(s, navy_bg, rule)

    if logo_path and Path(logo_path).exists():
        try:
            # White rounded plate behind the logo so the dark-navy wordmark
            # ("CONSTRUCTOR UNIVERSITY") stays readable on the navy header.
            from PIL import Image as _PILImage
            with _PILImage.open(str(logo_path)) as _im:
                _iw, _ih = _im.size
            logo_h = Inches(0.9)
            logo_w = int(logo_h * (_iw / _ih))
            logo_x, logo_y = Inches(0.5), Inches(0.4)
            pad_x, pad_y = Inches(0.18), Inches(0.12)
            plate = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                logo_x - pad_x, logo_y - pad_y,
                logo_w + 2 * pad_x, logo_h + 2 * pad_y,
            )
            plate.fill.solid()
            plate.fill.fore_color.rgb = WHITE
            plate.line.color.rgb = WHITE
            plate.shadow.inherit = False
            register_always_visible(s, plate)

            logo_shape = s.shapes.add_picture(str(logo_path), logo_x,
                                              logo_y, height=logo_h)
            register_always_visible(s, logo_shape)
            print(f"[info] Logo embedded on title slide ({Path(logo_path).name})")
        except Exception as e:
            print(f"[warn] could not embed logo: {e}")
    else:
        print("[warn] no logo file available for title slide")

    title = add_textbox(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(2.0),
                        "Empirical Evaluation and Cost Optimization\n"
                        "of Large Language Models in Azure Cloud Environments",
                        size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE)
    sub = add_textbox(s, Inches(0.7), Inches(3.7), Inches(11.9), Inches(0.5),
                      "Bachelor Thesis Defense, Computer Science",
                      size=18, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
    author = add_textbox(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.55),
                         "Matin Abaszada", size=24, bold=True, color=NAVY,
                         align=PP_ALIGN.CENTER)
    sup = add_textbox(s, Inches(1), Inches(5.55), Inches(11.3), Inches(0.45),
                      "Supervisor: Prof. Dr. Ivan Ovsyannikov", size=18,
                      color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    date = add_textbox(s, Inches(1), Inches(6.05), Inches(11.3), Inches(0.45),
                       "Constructor University, May 2026", size=16,
                       italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    register_always_visible(s, title, sub, author, sup, date)
    return s


def slide_motivation(prs):
    s = blank_slide(prs)
    add_header(s, "Why Cloud Hosting and Model Choice Both Matter",
               slide_num=2, section_label="MOTIVATION")

    block_w = Inches(4.0); block_h = Inches(3.0)
    top = Inches(1.4); gap = Inches(0.27); left0 = Inches(0.4)

    blocks = [
        ("The Expensive Reality",
         "Because of GPU clusters, cooling, dedicated operations and millions in CapEx, "
         "hosting LLMs locally is expensive."),
        ("The Shift to Cloud",
         "Enterprises move to cloud environments. One API call replaces an "
         "entire on-premise GPU server room. Pay only for what you use."),
        ("Azure as the Platform",
         "Microsoft Azure: a unified catalog of OpenAI / DeepSeek / Meta models, "
         "pay-per-token pricing, quotas, and rate limits."),
    ]
    for i, (title, body) in enumerate(blocks):
        x = left0 + i * (block_w + gap)
        shapes = add_card_block(s, x, top, block_w, block_h,
                                badge_num=i + 1, title=title, body=body,
                                body_size=14)
        register_group(s, *shapes)

    callout_shapes = add_callout(
        s, Inches(0.4), Inches(4.7), Inches(12.5), Inches(1.95),
        "In the cloud, “best model” ≠ “strongest model”.\n"
        "A top-accuracy model that costs ~20× more per request than the next-best "
        "may be the wrong enterprise choice.\n"
        "Accuracy, cost, and latency must all be evaluated together.",
        fill=AMBER_FILL, border=AMBER, size=16, bold=False)
    register_group(s, *callout_shapes)
    return s


def slide_research_gap(prs):
    s = blank_slide(prs)
    add_header(s, "Research Gap & Questions",
               slide_num=3, section_label="MOTIVATION")

    gap_title = add_textbox(s, Inches(0.5), Inches(1.3), Inches(6.5),
                            Inches(0.5), "Where the existing literature falls short",
                            size=22, bold=True, color=NAVY)
    gap_rule = add_filled_rect(s, Inches(0.5), Inches(1.85), Inches(1.2),
                               Inches(0.05), AZURE)
    gap_bullets = add_bullet_list(s, Inches(0.5), Inches(2.05), Inches(6.6),
                                  Inches(4.5), [
        "Most existing studies evaluate models in isolation: accuracy only, "
        "or cost only, never together.",
        "Optimization strategies (routing, cascading) are benchmarked against "
        "their own internal goal, not a unified deployment reward.",
        "Almost no studies test these strategies in a real cloud environment, "
        "where pricing, quotas, and rate limits directly shape latency and cost.",
    ], size=15)
    register_group(s, gap_title, gap_rule, gap_bullets)

    rx, rw = Inches(7.5), Inches(5.4)
    rq_panel = add_filled_rect(s, rx, Inches(1.3), rw, Inches(5.3), ICE_BLUE,
                               line_color=AZURE, line_width=Pt(1.0))
    rq_title = add_textbox(s, rx + Inches(0.3), Inches(1.45), rw - Inches(0.6),
                           Inches(0.6), "Research Questions",
                           size=22, bold=True, color=NAVY)
    rq_rule = add_filled_rect(s, rx + Inches(0.3), Inches(2.0),
                              Inches(1.2), Inches(0.04), AZURE)
    register_group(s, rq_panel, rq_title, rq_rule)

    rqs = [
        "How do Azure-hosted LLMs differ across accuracy, latency, and cost?",
        "Is reasoning-enabled inference economically justified?",
        "Can a proxy optimization layer "
        "improve cost-efficiency?",
    ]
    y = Inches(2.25)
    for i, rq in enumerate(rqs, 1):
        badge = add_filled_rect(s, rx + Inches(0.3), y, Inches(0.55),
                                Inches(0.55), AZURE)
        bt = s.shapes.add_textbox(rx + Inches(0.3), y, Inches(0.55), Inches(0.55))
        bt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = bt.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), str(i), size=18, bold=True, color=WHITE)
        rq_text = add_textbox(s, rx + Inches(1.05), y, rw - Inches(1.4),
                              Inches(1.4), rq, size=14, color=NEAR_BLACK)
        register_group(s, badge, bt, rq_text)
        y += Inches(1.4)
    return s


def slide_design(prs):
    s = blank_slide(prs)
    add_header(s, "Study Design at a Glance",
               slide_num=4, section_label="METHODOLOGY")

    rows = [
        ("6 Azure Models  (2 generations × 3 tiers)",
         "gpt-4.1-mini · gpt-5.4-mini    |    gpt-4.1 · gpt-5.4    |    "
         "o3-mini · gpt-5.4-pro"),
        ("5 Benchmark Datasets",
         "HumanEval (Python coding) · MBPP (Python coding) · "
         "MMLU-Pro (multi-subject multiple-choice) · "
         "GPQA (graduate-level science Q&A) · GSM8K (grade-school math word problems)"),
        ("Evaluated Under 4 Strategies",
         "Standalone  ·  Cascade  ·  Router  ·  Self-Consistency"),
    ]
    row_h = Inches(1.25); top = Inches(1.4); gap = Inches(0.25)
    left = Inches(1.0); width = Inches(11.3)

    prev_arrow = None
    for i, (title, body) in enumerate(rows):
        y = top + i * (row_h + gap)
        bg = add_filled_rect(s, left, y, width, row_h, ICE_BLUE,
                             line_color=AZURE, line_width=Pt(1.0))
        stripe = add_filled_rect(s, left, y, Inches(0.15), row_h, AZURE)
        ttl = add_textbox(s, left + Inches(0.4), y + Inches(0.12),
                          width - Inches(0.6), Inches(0.5),
                          title, size=18, bold=True, color=NAVY)
        bdy = add_textbox(s, left + Inches(0.4), y + Inches(0.62),
                          width - Inches(0.6), row_h - Inches(0.65),
                          body, size=15, color=NEAR_BLACK)
        if prev_arrow is not None:
            register_group(s, prev_arrow, bg, stripe, ttl, bdy)
        else:
            register_group(s, bg, stripe, ttl, bdy)
        if i < len(rows) - 1:
            prev_arrow = add_arrow_down(s, Inches(6.4),
                                        y + row_h + Inches(0.0),
                                        w=Inches(0.5), h=Inches(0.25))
        else:
            prev_arrow = None

    callout_shapes = add_callout(
        s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.7),
        "All strategies evaluated under the same economic reward function.",
        size=16, bold=True)
    register_group(s, *callout_shapes)
    return s


def slide_reward(prs):
    s = blank_slide(prs)
    add_header(s, "The Unified Evaluation Framework",
               slide_num=5, section_label="METHODOLOGY")

    eq_top = Inches(1.5)
    eq_bg = add_filled_rect(s, Inches(0.7), eq_top, Inches(11.9), Inches(1.4),
                            ICE_BLUE, line_color=AZURE, line_width=Pt(1.0))
    tb = s.shapes.add_textbox(Inches(0.7), eq_top + Inches(0.3),
                              Inches(11.9), Inches(0.9))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), "Reward  =  − ( ", size=28, bold=True, color=NEAR_BLACK)
    set_run(p.add_run(), "Cost", size=28, bold=True, color=GREEN)
    set_run(p.add_run(), "  +  ", size=28, bold=True, color=NEAR_BLACK)
    set_run(p.add_run(), "λ_latency × Latency", size=28, bold=True, color=AZURE)
    set_run(p.add_run(), "  +  ", size=28, bold=True, color=NEAR_BLACK)
    set_run(p.add_run(), "λ_error × Error", size=28, bold=True, color=ORANGE_RED)
    set_run(p.add_run(), " )", size=28, bold=True, color=NEAR_BLACK)
    register_group(s, eq_bg, tb)

    items = [
        ("$  Cost", "Cost of each request, deterministically calculated from tokens used × Azure price.", GREEN),
        ("⏱  Latency", "Wall-clock time per API call (seconds).", AZURE),
        ("✗  Error", "Represents whether the model response is correct or wrong, encoded as a binary value.", ORANGE_RED),
    ]
    block_w = Inches(3.95); gap = Inches(0.2); left0 = Inches(0.55); top = Inches(3.3)
    for i, (title, body, color) in enumerate(items):
        x = left0 + i * (block_w + gap)
        bg = add_filled_rect(s, x, top, block_w, Inches(1.45), WHITE,
                             line_color=color, line_width=Pt(1.5))
        ttl = add_textbox(s, x + Inches(0.2), top + Inches(0.12),
                          block_w - Inches(0.4), Inches(0.5),
                          title, size=20, bold=True, color=color)
        bdy = add_textbox(s, x + Inches(0.2), top + Inches(0.65),
                          block_w - Inches(0.4), Inches(0.75),
                          body, size=13, color=NEAR_BLACK)
        register_group(s, bg, ttl, bdy)

    callout_shapes = add_callout(
        s, Inches(0.55), Inches(5.0), Inches(12.25), Inches(2.0),
        "•  λ_latency = penalty per second of waiting.\n"
        "•  λ_error = penalty per wrong answer.\n"
        "•  Defaults: λ_latency = 0.01, λ_error = 1.0 (sensitivity swept).\n"
        "•  Higher reward (closer to zero) is better.",
        size=14, bold=True, align=PP_ALIGN.LEFT)
    register_group(s, *callout_shapes)
    return s


def slide_chart(prs, title, slide_num, section, chart_file,
                callouts=None, side_text=None, side_bullets=None,
                chart_left=None, chart_top=None, chart_w=None, chart_h=None):
    s = blank_slide(prs)
    add_header(s, title, slide_num=slide_num, section_label=section)

    if chart_left is None:
        cl, ct, cw, ch = Inches(0.6), Inches(1.3), Inches(8.6), Inches(5.5)
        if side_text is None and side_bullets is None and not callouts:
            cl, ct, cw, ch = Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.6)
    else:
        cl, ct, cw, ch = chart_left, chart_top, chart_w, chart_h

    chart_path = CHARTS / chart_file
    if chart_path.exists():
        pic = add_picture_fitted(s, chart_path, cl, ct, cw, ch)
        register_group(s, pic)
    else:
        miss = add_textbox(s, cl, ct, cw, ch,
                           f"[Chart missing: {chart_file}]",
                           size=18, color=ORANGE_RED, align=PP_ALIGN.CENTER,
                           anchor=MSO_ANCHOR.MIDDLE)
        register_group(s, miss)

    if side_text or side_bullets:
        sx, sw, sy = Inches(9.4), Inches(3.7), Inches(1.3)
        side_shapes = []
        if side_text:
            t = add_textbox(s, sx, sy, sw, Inches(0.5),
                            side_text["title"], size=16, bold=True, color=NAVY)
            side_shapes.append(t)
            sy += Inches(0.55)
        if side_bullets:
            b = add_bullet_list(s, sx, sy, sw, Inches(5.0), side_bullets,
                                size=13)
            side_shapes.append(b)
        register_group(s, *side_shapes)

    if callouts:
        for c in callouts:
            cs = add_callout(s, c["left"], c["top"], c["w"], c["h"], c["text"],
                             fill=c.get("fill", ICE_BLUE),
                             border=c.get("border", AZURE),
                             size=c.get("size", 13),
                             bold=c.get("bold", False))
            register_group(s, *cs)
    return s


def slide_cascade(prs):
    s = blank_slide(prs)
    add_header(s, "Cascade: Why not ask the small model first and only escalate when needed?",
               slide_num=10, section_label="OPTIMIZATION  •  1 / 3")

    lx, lw = Inches(0.4), Inches(4.5)
    head = add_textbox(s, lx, Inches(1.3), lw, Inches(0.45),
                       "How it works", size=16, bold=True, color=NAVY)
    register_group(s, head)

    by = Inches(1.85)
    nx, nw = lx + Inches(0.4), lw - Inches(0.8)

    n1 = add_diag_node(s, nx, by, nw, Inches(0.55), "Query")
    register_group(s, *n1)

    arr1 = add_arrow_down(s, lx + lw / 2 - Inches(0.15), by + Inches(0.6))
    n2 = add_diag_node(s, nx, by + Inches(0.85), nw, Inches(0.55),
                       "Small Model (cheap)")
    register_group(s, arr1, *n2)

    arr2 = add_arrow_down(s, lx + lw / 2 - Inches(0.15), by + Inches(1.45))
    n3 = add_diag_node(s, nx, by + Inches(1.7), nw, Inches(0.55),
                       "Confidence ≥ Threshold ?", color=AMBER, fill=AMBER_FILL)
    register_group(s, arr2, *n3)

    yes_t = add_textbox(s, lx + Inches(0.1), by + Inches(2.4), lw / 2,
                        Inches(0.4), "YES  →  return", size=12, bold=True,
                        color=GREEN, align=PP_ALIGN.CENTER)
    no_t = add_textbox(s, lx + Inches(2.2), by + Inches(2.4), lw / 2,
                       Inches(0.4), "NO  →  Large Model", size=12, bold=True,
                       color=ORANGE_RED, align=PP_ALIGN.CENTER)
    n4 = add_diag_node(s, nx, by + Inches(2.85), nw, Inches(0.55),
                       "Large Model (final answer)", color=ORANGE_RED)
    register_group(s, yes_t, no_t, *n4)

    # --- Table 7: Cascade reward per configuration with standalone baselines
    # Header bar ends at y=1.01"; bottom amber callout starts at y=6.15".
    # Available vertical span = 5.14".  Center a (table+caption) block of
    # ~3.05" height inside it: top = 1.01 + (5.14 - 3.05)/2 ~= 2.06"
    tbl_left, tbl_top = Inches(5.45), Inches(2.05)
    tbl_w = Inches(7.45)
    headers = ["Configuration", "T=60", "T=75", "T=90",
               "Small alone", "Large alone"]
    rows = [
        ("gpt-4.1-mini → gpt-4.1", "−0.601", "−0.602", "−0.606",
         "−0.582", "−0.527"),
        ("gpt-4.1-mini → gpt-5.4", "−0.601", "−0.601", "−0.601",
         "−0.582", "−0.444"),
        ("gpt-5.4-mini → gpt-4.1", "−0.680", "−0.689", "−0.708",
         "−0.569", "−0.527"),
        ("gpt-5.4-mini → gpt-5.4", "−0.680", "−0.680", "−0.681",
         "−0.569", "−0.444"),
    ]
    # Best (least-negative reward = lowest penalty) cell in each data row.
    # All four rows have "Large alone" (column 5) as the maximum.
    best_cells = {(ri + 1, 5) for ri in range(len(rows))}

    n_rows, n_cols = len(rows) + 1, len(headers)
    # Row heights: header slightly larger; data rows uniform.
    row_heights = [Inches(0.45)] + [Inches(0.50)] * len(rows)
    tbl_h = sum(row_heights, Inches(0))
    table_shape = s.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top,
                                     tbl_w, tbl_h)
    table = table_shape.table

    # --- Strip default table style so our cell fills/borders are not overridden
    tbl_el = table._tbl
    tblPr = tbl_el.find(qn("a:tblPr"))
    if tblPr is not None:
        tblPr.set("firstRow", "0")
        tblPr.set("bandRow", "0")
        tblPr.set("firstCol", "0")
        for tsid in tblPr.findall(qn("a:tableStyleId")):
            tblPr.remove(tsid)

    col_w = [Inches(2.30), Inches(0.78), Inches(0.78), Inches(0.78),
             Inches(1.40), Inches(1.41)]
    for ci, w in enumerate(col_w):
        try:
            table.columns[ci].width = w
        except Exception:
            pass
    for ri, h in enumerate(row_heights):
        try:
            table.rows[ri].height = h
        except Exception:
            pass

    def _set_cell_borders(cell, color_hex, w_emu="12700"):
        """Set all four borders of a table cell to a solid color (RGB hex)."""
        tcPr = cell._tc.get_or_add_tcPr()
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            for el in tcPr.findall(qn(tag)):
                tcPr.remove(el)
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            ln = etree.SubElement(tcPr, qn(tag))
            ln.set("w", w_emu)
            ln.set("cap", "flat")
            ln.set("cmpd", "sng")
            ln.set("algn", "ctr")
            fill = etree.SubElement(ln, qn("a:solidFill"))
            clr = etree.SubElement(fill, qn("a:srgbClr"))
            clr.set("val", color_hex)
            etree.SubElement(ln, qn("a:prstDash")).set("val", "solid")
            etree.SubElement(ln, qn("a:round"))

    border_hex = "{:02X}{:02X}{:02X}".format(NAVY[0], NAVY[1], NAVY[2])

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.margin_left = Inches(0.03); cell.margin_right = Inches(0.03)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_cell_borders(cell, border_hex)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = ""
        set_run(p.add_run(), h, size=11, bold=True, color=WHITE)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.margin_left = Inches(0.03); cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if (ri + 1, ci) in best_cells:
                cell.fill.fore_color.rgb = GREEN_FILL
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else ALT_ROW
            _set_cell_borders(cell, border_hex)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = ""
            bold = (ci == 0) or ((ri + 1, ci) in best_cells)
            color = GREEN if (ri + 1, ci) in best_cells else NEAR_BLACK
            set_run(p.add_run(), val, size=11, bold=bold, color=color)
    register_group(s, table_shape)

    # Caption directly below the table.
    cap_top = tbl_top + tbl_h + Inches(0.06)
    cap = add_textbox(s, tbl_left, cap_top, tbl_w, Inches(0.3),
                      "Table 7: Cascade reward at default weights "
                      "(λ_l = 0.01, λ_e = 1.0).  No cascade beats its "
                      "large-alone baseline.",
                      size=10, italic=True, color=GRAY,
                      align=PP_ALIGN.CENTER)
    register_group(s, cap)

    callout_shapes = add_callout(
        s, Inches(0.6), Inches(6.15), Inches(12.4), Inches(0.85),
        "Maximum escalation rate: only 7.3%, even at T = 90.  "
        "Modern LLMs are overconfident; the gate almost never opens.",
        fill=AMBER_FILL, border=AMBER, size=14, bold=True)
    register_group(s, *callout_shapes)
    return s


def slide_self_consistency(prs):
    s = blank_slide(prs)
    add_header(s, "Self-Consistency: Why not let the small model double-check its own answer?",
               slide_num=11, section_label="OPTIMIZATION  •  2 / 3")

    lx, lw = Inches(0.4), Inches(4.4)
    head = add_textbox(s, lx, Inches(1.3), lw, Inches(0.45),
                       "How it works", size=16, bold=True, color=NAVY)
    register_group(s, head)

    by = Inches(1.85)
    nx, nw = lx + Inches(0.4), lw - Inches(0.8)
    q = add_diag_node(s, nx, by, nw, Inches(0.55), "Query")
    register_group(s, *q)

    py = by + Inches(0.95)
    parallel_shapes = []
    run_arrows = []
    for i in range(3):
        x = lx + Inches(0.3) + i * Inches(1.25)
        run_arrows.append(
            add_arrow_down(s, x + Inches(0.5), py - Inches(0.3),
                           w=Inches(0.15), h=Inches(0.25))
        )
        node = add_diag_node(s, x, py, Inches(1.15), Inches(0.5),
                             f"Run {i+1}", size=12)
        parallel_shapes.extend(node)
    register_group(s, *run_arrows, *parallel_shapes)

    arr = add_arrow_down(s, lx + lw / 2 - Inches(0.15), py + Inches(0.6))
    mv_y = py + Inches(0.95)
    mv = add_diag_node(s, nx, mv_y, nw, Inches(0.55),
                       "Majority Vote → Final Answer",
                       color=GREEN, fill=GREEN_FILL)
    register_group(s, arr, *mv)

    cs = add_callout(s, lx, Inches(4.7), lw + Inches(0.05), Inches(2.0),
                     "gpt-4.1-mini N=3:  +6.8 pp accuracy.\n"
                     "MMLU-Pro: 27.8% → 44.0%.\n\n"
                     "Modest gain: cost = exactly 3×,\n"
                     "and accuracy stays below the next\n"
                     "standalone tier (68.3%).",
                     size=13, bold=False)
    register_group(s, *cs)

    pic = add_picture_fitted(s, CHARTS / "chartSC5_selfcons_accuracy_comparison.png",
                             Inches(5.0), Inches(1.25), Inches(8.1), Inches(5.4))
    register_group(s, pic)
    return s


def slide_router(prs):
    s = blank_slide(prs)
    add_header(s, "Router: Why not let an intelligent LLM decide which model to use?",
               slide_num=12, section_label="OPTIMIZATION  •  3 / 3")

    lx, lw = Inches(0.4), Inches(4.5)
    head = add_textbox(s, lx, Inches(1.3), lw, Inches(0.45),
                       "How it works", size=16, bold=True, color=NAVY)
    register_group(s, head)

    by = Inches(1.85)
    nx, nw = lx + Inches(0.4), lw - Inches(0.8)
    q = add_diag_node(s, nx, by, nw, Inches(0.55), "Query")
    register_group(s, *q)

    arr1 = add_arrow_down(s, lx + lw / 2 - Inches(0.15), by + Inches(0.6))
    rt = add_diag_node(s, nx, by + Inches(0.85), nw, Inches(0.55),
                       "Router  (gpt-5.4-mini)", color=AMBER, fill=AMBER_FILL)
    register_group(s, arr1, *rt)

    arr2 = add_arrow_down(s, lx + lw / 2 - Inches(0.15), by + Inches(1.45))
    bb = by + Inches(1.7)
    easy = add_diag_node(s, lx + Inches(0.3), bb, Inches(1.95), Inches(0.6),
                         "EASY → Small\n(gpt-4.1-mini)", color=GREEN,
                         fill=GREEN_FILL, size=11)
    hard = add_diag_node(s, lx + Inches(2.4), bb, Inches(1.95), Inches(0.6),
                         "HARD → Large\n(gpt-4.1)", color=ORANGE_RED,
                         fill=RED_FILL, size=11)
    register_group(s, arr2, *easy, *hard)

    cs = add_callout(s, lx, Inches(4.55), lw + Inches(0.05), Inches(2.05),
                     "Accuracy clearly above the small baseline, at a fraction\n"
                     "of the cost of the large model.\n\n"
                     "Trade-off: latency stays close to the large model,\n"
                     "because every query first passes through the router.",
                     fill=GREEN_FILL, border=GREEN, size=12, bold=False)
    register_group(s, *cs)

    pic = add_picture_fitted(s, CHARTS / "chartR3_router_best_config_heatmap.png",
                             Inches(5.0), Inches(1.25), Inches(8.1), Inches(5.4))
    register_group(s, pic)
    return s


def slide_summary_table(prs):
    s = blank_slide(prs)
    add_header(s, "Head-to-Head: Best Configuration per Strategy",
               slide_num=13, section_label="RESULTS")

    headers = ["Strategy", "Best Configuration", "Accuracy", "Latency",
               "Cost (m$)", "Reward"]
    rows = [
        ("Standalone", "gpt-5.4", "90.8%", "34.7 s", "5.370", "−0.444"),
        ("Router", "rtr=gpt-5.4-mini → gpt-4.1-mini / gpt-4.1",
         "73.1%", "22.9 s", "0.344", "−0.498"),
        ("Self-Consist.", "gpt-4.1-mini  N=3", "64.4%", "18.6 s", "0.765", "−0.543"),
        ("Cascade", "gpt-4.1-mini → gpt-5.4  T=75",
         "57.7%", "17.7 s", "0.261", "−0.601"),
    ]
    highlights = {
        # Accuracy (col 2): best Standalone 90.8%, worst Cascade 57.7%
        (0, 2): GREEN, (3, 2): ORANGE_RED,
        # Latency (col 3): best Cascade 17.7s, worst Standalone 34.7s
        (3, 3): GREEN, (0, 3): ORANGE_RED,
        # Cost (col 4): best Cascade 0.261, worst Standalone 5.370
        (3, 4): GREEN, (0, 4): ORANGE_RED,
        # Reward (col 5): best Standalone -0.444, worst Cascade -0.601
        (0, 5): GREEN, (3, 5): ORANGE_RED,
    }

    n_rows = len(rows) + 1; n_cols = len(headers)
    left, top = Inches(0.5), Inches(1.3)
    width, height = Inches(12.3), Inches(3.6)

    table_shape = s.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Strip default table style so cell borders/fills are not overridden
    tbl_el = table._tbl
    tblPr = tbl_el.find(qn("a:tblPr"))
    if tblPr is not None:
        tblPr.set("firstRow", "0")
        tblPr.set("bandRow", "0")
        tblPr.set("firstCol", "0")
        for tsid in tblPr.findall(qn("a:tableStyleId")):
            tblPr.remove(tsid)

    border_hex = "{:02X}{:02X}{:02X}".format(NAVY[0], NAVY[1], NAVY[2])

    def _set_cell_borders(cell, color_hex, w_emu="12700"):
        tcPr = cell._tc.get_or_add_tcPr()
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            for el in tcPr.findall(qn(tag)):
                tcPr.remove(el)
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            ln = etree.SubElement(tcPr, qn(tag))
            ln.set("w", w_emu); ln.set("cap", "flat")
            ln.set("cmpd", "sng"); ln.set("algn", "ctr")
            fill = etree.SubElement(ln, qn("a:solidFill"))
            clr = etree.SubElement(fill, qn("a:srgbClr"))
            clr.set("val", color_hex)
            etree.SubElement(ln, qn("a:prstDash")).set("val", "solid")
            etree.SubElement(ln, qn("a:round"))

    col_w = [Inches(1.5), Inches(4.4), Inches(1.4), Inches(1.4),
             Inches(1.6), Inches(1.6)]
    for ci, w in enumerate(col_w):
        try:
            table.columns[ci].width = w
        except Exception:
            pass

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_cell_borders(cell, border_hex)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = ""
        set_run(p.add_run(), h, size=14, bold=True, color=WHITE)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            highlight = highlights.get((ri, ci))
            cell.fill.solid()
            if highlight == GREEN:
                cell.fill.fore_color.rgb = GREEN_FILL
            elif highlight == ORANGE_RED:
                cell.fill.fore_color.rgb = RED_FILL
            elif highlight is not None:
                cell.fill.fore_color.rgb = ICE_BLUE
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else ALT_ROW
            _set_cell_borders(cell, border_hex)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci != 1 else PP_ALIGN.LEFT
            p.text = ""
            bold = (ci == 0) or (ci == 5)
            set_run(p.add_run(), val, size=13, bold=bold, color=NEAR_BLACK)

    register_group(s, table_shape)

    note = add_textbox(s, Inches(0.5), Inches(4.95), Inches(12.3),
                       Inches(0.35),
                       "Cost (m$) = milli-dollars per request "
                       "(1 m$ = 0.001 USD = one thousandth of a US dollar). "
                       "Computed as token counts × Azure list price.",
                       size=11, italic=True, color=GRAY,
                       align=PP_ALIGN.CENTER)
    register_group(s, note)

    cs = add_callout(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4),
                     "•  At default weights, the best reward belongs to a standalone model.\n"
                     "•  This ordering changes as the latency and error penalty weights are varied, as the next slide shows.",
                     size=15, bold=True)
    register_group(s, *cs)
    return s


def slide_unified_heatmap(prs):
    s = blank_slide(prs)
    add_header(s, "Best Model & Optimization Strategy Across the Penalty Plane",
               slide_num=14, section_label="RESULTS")

    pic = add_picture_fitted(s, CHARTS / "chartUNI_unified_best_strategy_heatmap_router_diagonal.png",
                             Inches(0.4), Inches(1.25),
                             Inches(8.9), Inches(5.7))
    register_group(s, pic)

    sx, sy, sw = Inches(9.4), Inches(1.3), Inches(3.7)
    head = add_textbox(s, sx, sy, sw, Inches(0.5),
                       "Who wins each cell",
                       size=16, bold=True, color=NAVY)
    rule = add_filled_rect(s, sx, sy + Inches(0.55),
                           Inches(1.2), Inches(0.04), AZURE)
    bullets = add_bullet_list(s, sx, sy + Inches(0.7), sw, Inches(3.3), [
        "SC gpt-4.1-mini (N=3): removes noise from a cheap and fast model, but accuracy still trails the standalone models",
        "gpt-5.4: slower than SC, but reaches higher accuracy (fewer errors) — ★ default sits here",
        "gpt-5.4-pro: highest accuracy of all, wins where error penalty is high and latency penalty is low",
        "Router: accurate at a much lower cost than the large models, but latency stays high",
    ], size=13)
    register_group(s, head, rule, bullets)
    return s


def slide_conclusions(prs):
    s = blank_slide(prs)
    add_header(s, "Conclusions: Answering the Research Questions",
               slide_num=15, section_label="CONCLUSIONS")

    items = [
        ("RQ 1",
         "How do Azure LLMs differ in accuracy, latency, and cost?",
         "Two-tier structure: 15 pp accuracy gap and ≈7× cost gap at the tier boundary "
         "(o3-mini vs gpt-4.1).  No model Pareto-dominates.  gpt-5.4 = best balanced choice at default weights.",
         AZURE),
        ("RQ 2",
         "Is reasoning-enabled inference economically justified?",
         "Only for accuracy-critical, latency-tolerant tasks.  gpt-5.4-pro adds "
         "+2.4 pp at 20× the cost and 1.5× the latency. Not justified for general workloads.",
         AMBER),
        ("RQ 3",
         "Can proxy optimization improve cost-efficiency?",
         "•  Cascade fails to optimize because neutralized by LLM overconfidence.\n"
         "•  Self-Consistency optimizes the cheap/fast regime by reducing noise.\n"
         "•  Router optimizes the accuracy-vs-cost trade-off, winning a band of the decision map.",
         GREEN),
    ]

    top = Inches(1.3); h = Inches(1.7); gap = Inches(0.18)
    for i, (tag, q, a, color) in enumerate(items):
        y = top + i * (h + gap)
        tag_bg = add_filled_rect(s, Inches(0.5), y, Inches(1.3), h, color)
        tb = s.shapes.add_textbox(Inches(0.5), y, Inches(1.3), h)
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), tag, size=24, bold=True, color=WHITE)
        body_bg = add_filled_rect(s, Inches(1.85), y, Inches(11.0), h,
                                  WHITE, line_color=color, line_width=Pt(1.25))
        q_tb = add_textbox(s, Inches(2.05), y + Inches(0.12),
                           Inches(10.7), Inches(0.55), q,
                           size=15, bold=True, color=NAVY)
        a_tb = s.shapes.add_textbox(Inches(2.05), y + Inches(0.65),
                                    Inches(10.7), h - Inches(0.7))
        a_tf = a_tb.text_frame
        a_tf.word_wrap = True
        a_tf.vertical_anchor = MSO_ANCHOR.TOP
        for li, line in enumerate(a.split("\n")):
            p_a = a_tf.paragraphs[0] if li == 0 else a_tf.add_paragraph()
            p_a.alignment = PP_ALIGN.LEFT
            set_run(p_a.add_run(), line, size=14, color=NEAR_BLACK)
        register_group(s, tag_bg, tb, body_bg, q_tb, a_tb)
    return s


def slide_limits(prs):
    s = blank_slide(prs)
    add_header(s, "Limitations & Future Work",
               slide_num=17, section_label="CONCLUSIONS")

    col_w = Inches(6.0); top = Inches(1.3)

    lim_head = add_filled_rect(s, Inches(0.5), top, col_w, Inches(0.6), NAVY)
    lim_head_t = add_textbox(s, Inches(0.5), top, col_w, Inches(0.6),
                             "Limitations", size=20, bold=True, color=WHITE,
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lim_panel = add_filled_rect(s, Inches(0.5), top + Inches(0.6), col_w,
                                Inches(3.8), ICE_BLUE,
                                line_color=AZURE, line_width=Pt(1.0))
    lim_bul = add_bullet_list(s, Inches(0.7), top + Inches(0.75),
                              col_w - Inches(0.4), Inches(3.7), [
        "6 models from gpt-4.1 / 5.4 / o3 families only. Other Azure models not covered.",
        "5 benchmark families; no long-context or multi-turn dialogue.",
        "Latency measured under one quota / region; may vary.",
        "Cascade used raw confidence without post-hoc calibration.",
    ], size=14)
    register_group(s, lim_head, lim_head_t, lim_panel, lim_bul)

    fx = Inches(6.85)
    fw_head = add_filled_rect(s, fx, top, col_w, Inches(0.6), AZURE)
    fw_head_t = add_textbox(s, fx, top, col_w, Inches(0.6),
                            "Future Work", size=20, bold=True, color=WHITE,
                            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    fw_panel = add_filled_rect(s, fx, top + Inches(0.6), col_w, Inches(3.8),
                               GREEN_FILL, line_color=GREEN, line_width=Pt(1.0))
    fw_bul = add_bullet_list(s, fx + Inches(0.2), top + Inches(0.75),
                             col_w - Inches(0.4), Inches(3.7), [
        "Calibrated cascade (post-hoc / verifier-based escalation).",
        "Lightweight dedicated router models; hybrid strategies.",
        "Extended benchmarks: domain-specific, long-context, dialogue.",
        "Online deployment framework with real-time quota signals.",
    ], size=14)
    register_group(s, fw_head, fw_head_t, fw_panel, fw_bul)

    ty_bg = add_filled_rect(s, Inches(0.5), Inches(5.5), Inches(12.35),
                            Inches(1.3), NAVY)
    ty_t = add_textbox(s, Inches(0.5), Inches(5.5), Inches(12.35), Inches(0.7),
                       "Thank you!", size=32, bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ty_q = add_textbox(s, Inches(0.5), Inches(6.15), Inches(12.35),
                       Inches(0.55), "Questions welcome.",
                       size=18, italic=True, color=ICE_BLUE,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    register_group(s, ty_bg, ty_t, ty_q)
    return s


def slide_thanks(prs):
    s = blank_slide(prs)
    navy_bg = add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    register_always_visible(s, navy_bg)

    thanks = add_textbox(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(1.6),
                         "Thank You",
                         size=72, bold=True, color=WHITE,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    sub = add_textbox(s, Inches(0.7), Inches(3.7), Inches(11.9), Inches(0.7),
                      "Questions & Discussion",
                      size=24, italic=True, color=WHITE,
                      align=PP_ALIGN.CENTER)
    rule = add_filled_rect(s, Inches(4.5), Inches(4.55), Inches(4.333),
                           Inches(0.05), AZURE)
    title = add_textbox(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.0),
                        "Empirical Evaluation and Cost Optimization\n"
                        "of Large Language Models in Azure Cloud Environments",
                        size=20, bold=True, color=WHITE,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    author = add_textbox(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.5),
                         "Matin Abaszada", size=20, bold=True, color=WHITE,
                         align=PP_ALIGN.CENTER)
    sup = add_textbox(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.45),
                      "Supervisor: Prof. Dr. Ivan Ovsyannikov   •   "
                      "Constructor University, May 2026",
                      size=14, italic=True, color=WHITE,
                      align=PP_ALIGN.CENTER)
    register_always_visible(s, thanks, sub, rule, title, author, sup)
    return s


def slide_section_break(prs, label, subtitle):
    s = blank_slide(prs)
    bg = add_filled_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rule = add_filled_rect(s, Inches(1.5), Inches(3.6), Inches(2.5),
                           Inches(0.07), AZURE)
    lab = add_textbox(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(0.7),
                      label.upper(), size=18, bold=True, color=AZURE,
                      align=PP_ALIGN.LEFT)
    sub = add_textbox(s, Inches(1.0), Inches(3.85), Inches(11.3), Inches(2.0),
                      subtitle, size=36, bold=True, color=WHITE,
                      align=PP_ALIGN.LEFT)
    register_always_visible(s, bg, rule, lab, sub)
    return s


def slide_backup_chart(prs, title, num, chart_file, key_insight):
    s = blank_slide(prs)
    add_header(s, title, slide_num=None, section_label=f"BACKUP  •  B{num}")

    pic = add_picture_fitted(s, CHARTS / chart_file,
                             Inches(0.5), Inches(1.3), Inches(8.7), Inches(5.7))
    register_group(s, pic)

    cs = add_callout(s, Inches(9.4), Inches(1.3), Inches(3.7), Inches(5.6),
                     key_insight, size=13, bold=False)
    register_group(s, *cs)
    return s


def slide_backup_text(prs, title, num, body_blocks):
    s = blank_slide(prs)
    add_header(s, title, slide_num=None, section_label=f"BACKUP  •  B{num}")
    y = Inches(1.2)
    line_h = 0.22  # inches per wrapped body line at 13pt
    chars_per_line = 115
    for tag, txt in body_blocks:
        head = add_textbox(s, Inches(0.6), y, Inches(12.1), Inches(0.32),
                           tag, size=16, bold=True, color=NAVY)
        y += Inches(0.34)
        rule = add_filled_rect(s, Inches(0.6), y, Inches(12.1),
                               Inches(0.04), AZURE)
        y += Inches(0.10)
        n_lines = max(1, math.ceil(len(txt) / chars_per_line))
        body_h = line_h * n_lines + 0.05
        body_tb = add_textbox(s, Inches(0.6), y, Inches(12.1), Inches(body_h),
                              txt, size=13, color=NEAR_BLACK)
        y += Inches(body_h + 0.18)
        register_group(s, head, rule, body_tb)
    return s


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def main():
    print("[info] Building presentation ...")
    logo = extract_logo()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = []
    slides.append(slide_title(prs, logo))                         # 1
    slides.append(slide_motivation(prs))                          # 2
    slides.append(slide_research_gap(prs))                        # 3
    slides.append(slide_design(prs))                              # 4
    slides.append(slide_reward(prs))                              # 5

    slides.append(slide_chart(
        prs, "Accuracy: A Two-Tier Structure", 6, "RESULTS  •  STANDALONE",
        "chart1_accuracy.png",
        side_text={"title": "Two clear tiers"},
        side_bullets=[
            "Top tier (>80%): gpt-5.4-pro 93.2%, gpt-5.4 90.8%, o3-mini 83.3%",
            "Bottom tier (<70%): gpt-4.1 68.3%, gpt-5.4-mini 60.6%, gpt-4.1-mini 57.6%",
            "Gap between tiers ≈ 15 pp, larger than the spread within either tier",
        ],
        chart_left=Inches(0.4), chart_top=Inches(1.3),
        chart_w=Inches(8.9), chart_h=Inches(5.7),
        callouts=[{
            "left": Inches(9.4), "top": Inches(5.5),
            "w": Inches(3.7), "h": Inches(1.4),
            "text": "Not gradual improvement; a structural jump.",
            "size": 14, "bold": True,
        }],
    ))                                                            # 6

    slides.append(slide_chart(
        prs, "Cost vs. Accuracy: ~20× Gap at the Pareto Frontier", 7, "RESULTS  •  STANDALONE",
        "chart3_cost_vs_accuracy.png",
        side_text={"title": "Pareto frontier"},
        side_bullets=[
            "gpt-5.4: 90.8% accuracy at 23× lower cost than gpt-5.4-pro",
            "gpt-5.4-pro: +2.4 pp over gpt-5.4 at ~20× the cost",
            "o3-mini and gpt-4.1 are Pareto-dominated",
        ],
        chart_left=Inches(0.4), chart_top=Inches(1.3),
        chart_w=Inches(8.9), chart_h=Inches(5.7),
        callouts=[{
            "left": Inches(9.4), "top": Inches(5.5),
            "w": Inches(3.7), "h": Inches(1.4),
            "text": "Cost spans almost three orders of magnitude.",
            "size": 13, "bold": True,
        }],
    ))                                                            # 7

    slides.append(slide_chart(
        prs, "Economic Reward: The Ranking Changes", 8,
        "RESULTS  •  STANDALONE", "chart4_reward.png",
        side_text={"title": "Most accurate ≠ best reward"},
        side_bullets=[
            "gpt-5.4-pro: #1 accuracy → #5 reward (cost + latency penalty)",
            "o3-mini: #3 accuracy → #6 reward (65.7 s mean latency)",
            "gpt-5.4: #2 accuracy → #1 reward (best balance)",
        ],
        chart_left=Inches(0.4), chart_top=Inches(1.3),
        chart_w=Inches(8.9), chart_h=Inches(5.7),
    ))                                                            # 8

    slides.append(slide_chart(
        prs, "Which Model Wins, and When?", 9, "RESULTS  •  STANDALONE",
        "chart7_best_model_heatmap.png",
        side_bullets=[
            "Each cell = winning model at those (λ_error, λ_latency) weights",
            "★  marks the default deployment setting",
            "Mini wins under high latency penalty; gpt-5.4 wins the centre; "
            "gpt-5.4-pro wins under high error penalty",
        ],
        chart_left=Inches(0.4), chart_top=Inches(1.3),
        chart_w=Inches(8.9), chart_h=Inches(5.7),
    ))                                                            # 9

    slides.append(slide_cascade(prs))                             # 10
    slides.append(slide_self_consistency(prs))                    # 11
    slides.append(slide_router(prs))                              # 12
    slides.append(slide_summary_table(prs))                       # 13
    slides.append(slide_unified_heatmap(prs))                     # 14
    slides.append(slide_conclusions(prs))                         # 15
    slides.append(slide_thanks(prs))                              # 16

    slides.append(slide_section_break(
        prs, "Backup", "Additional charts & detail for Q&A"))

    slides.append(slide_backup_chart(
        prs, "Per-Dataset Accuracy Breakdown", 1,
        "chart2_accuracy_per_dataset.png",
        "Coding benchmarks are saturated (80–100%).  GPQA and MMLU-Pro "
        "are the most discriminating datasets."))

    slides.append(slide_backup_chart(
        prs, "Self-Consistency Decision Map", 2,
        "chartSC3_selfcons_best_config_heatmap.png",
        "gpt-4.1-mini N=3 dominates the upper-left "
        "(latency-tolerant, accuracy-important) region."))

    slides.append(slide_backup_chart(
        prs, "Router Detail: Routing Behavior by Dataset", 3,
        "chartR1_router_overview.png",
        "Coding tasks are almost never escalated; GPQA gets 40%+ "
        "escalation rate."))

    slides.append(slide_backup_chart(
        prs, "Per-Dataset Decision Maps", 4,
        "chartUNI_unified_best_strategy_heatmap_per_dataset.png",
        "Winner per (λ_error, λ_latency) cell, broken down by dataset. "
        "Shows that the best strategy is not uniform across benchmarks — "
        "e.g. coding tasks favor the small model almost everywhere, while "
        "GPQA and MMLU-Pro give the Router and the large models more room."))

    slides.append(slide_backup_text(
        prs, "Reward Formula: Full Detail", 4,
        [("Equation",
          "Reward = − ( Cost  +  λ_latency × Latency  +  λ_error × Error )"),
         ("Why a negative penalty?",
          "A positive utility would require an arbitrary ceiling for what a "
          "“perfect” answer is worth. The negative-penalty form has a natural "
          "zero (no cost, no latency, no error), which is the theoretical "
          "ideal. Every real model can only be penalized from there."),
         ("Units & defaults",
          "Cost in dollars per request (token counts × Azure price). Latency "
          "in seconds. Error binary (0 or 1). Default weights: "
          "λ_latency = 0.01, λ_error = 1.0. Sensitivity is swept on a 2-D grid."),
         ("Interpretation",
          "−0.444 means the average request incurs 0.444 units of combined "
          "penalty.  Higher (closer to 0) = better.")]))

    print("[info] Adding grouped entrance animations ...")
    for i, sl in enumerate(slides):
        try:
            add_grouped_animations(sl)
        except Exception as e:
            print(f"[warn] animation skipped on slide {i+1}: {e}")

    prs.save(str(OUT_PATH))
    print(f"[done] Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
